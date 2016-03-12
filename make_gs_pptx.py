#!/usr/bin/env python
from pptx import Presentation
from pptx.util import Inches, Pt

Titles_And_Blurbs_Width = Inches(7)
Title_Height = Inches(0.5)


def AddBill(slide, Title, Committee_Name, Committee_Recommendation, Bill_Type, Bill_Summary, Blurbs):

  Title_Box = slide.shapes.add_textbox(0,0, Titles_And_Blurbs_Width, Title_Height)
  tf = Title_Box.text_frame
  fill = tf.fill
  fill.solid()
  fill.fore_color.rgb = RGBColor(0x01, 0x23, 0x45)
  fill.fore_color.brightness = 0.25
  fill.background()
  tf.fill = fill

  tf.text = Title



prs = Presentation()
temp = prs.slide_width
prs.slide_width = prs.slide_height
prs.slide_height = temp

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
Blurbs=['Does some stuff', 'does some othter stuff', 'should do that']

AddBill(slide, "This is a title", "Committe", "ITL", "Anti", "Does stuff", Blurbs)

#title = slide.shapes.title
#subtitle = slide.placeholders[1]


prs.save('test.pptx')
